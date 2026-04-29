"""
Build presentation/slides.pptx straight from the outline.

Run once:
    pip install python-pptx==0.6.23
    python build_pptx.py

Then open slides.pptx in PowerPoint to embed the demo video on the
"Live Demo" slide (Insert > Video > This Device > utd_demo.mp4 >
Set Start = Automatically). The narrative figures are already inserted.
Speaker notes are pre-filled per slide.

The narrative this deck follows is documented in
`presentation/NARRATIVE.md`.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

HERE     = os.path.dirname(os.path.abspath(__file__))
ROOT     = os.path.dirname(HERE)
OUT      = os.path.join(HERE, "slides.pptx")
FIGURES  = os.path.join(ROOT, "report", "figures")

UTD_ORANGE = RGBColor(0xC7, 0x55, 0x00)
UTD_GREEN  = RGBColor(0x15, 0x4E, 0x37)
DARK_TEXT  = RGBColor(0x22, 0x22, 0x22)
GREY_TEXT  = RGBColor(0x66, 0x66, 0x66)


def _style_title(shape, color=UTD_GREEN, size=28):
    for run in shape.text_frame.paragraphs[0].runs:
        run.font.color.rgb = color
        run.font.size = Pt(size)
        run.font.bold = True


def _style_body(text_frame, size=18):
    for p in text_frame.paragraphs:
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = DARK_TEXT


def add_title_slide(prs, title, subtitle, authors):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = f"{subtitle}\n{authors}"
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = UTD_ORANGE
        run.font.size = Pt(40)
        run.font.bold = True
    return slide


def add_bullet_slide(prs, title, bullets, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title, size=30)
    body = slide.placeholders[1].text_frame
    body.word_wrap = True
    body.text = bullets[0]
    for b in bullets[1:]:
        p = body.add_paragraph()
        p.text = b
        p.level = 0
    _style_body(body, size=20)
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_image_slide(prs, title, image_path, bullets=None, notes="",
                    img_left=1.0, img_top=1.4, img_width=11.0):
    """Title + a hero image, optionally with bullets stacked below."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title
    _style_title(slide.shapes.title, size=28)

    if os.path.exists(image_path):
        slide.shapes.add_picture(
            image_path, Inches(img_left), Inches(img_top),
            width=Inches(img_width),
        )

    if bullets:
        tbox = slide.shapes.add_textbox(Inches(0.7), Inches(6.3),
                                        Inches(12), Inches(1))
        tf = tbox.text_frame
        tf.word_wrap = True
        tf.text = bullets[0]
        for b in bullets[1:]:
            p = tf.add_paragraph()
            p.text = b
        _style_body(tf, size=14)

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_two_image_slide(prs, title, left_image, right_image,
                        left_caption, right_caption, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title, size=28)

    img_w = 6.0
    if os.path.exists(left_image):
        slide.shapes.add_picture(left_image, Inches(0.4), Inches(1.5),
                                 width=Inches(img_w))
    if os.path.exists(right_image):
        slide.shapes.add_picture(right_image, Inches(6.9), Inches(1.5),
                                 width=Inches(img_w))

    for x, caption in ((0.4, left_caption), (6.9, right_caption)):
        cap = slide.shapes.add_textbox(Inches(x), Inches(5.5),
                                       Inches(img_w), Inches(0.6))
        cap.text_frame.text = caption
        for p in cap.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(14)
                run.font.color.rgb = GREY_TEXT
                run.font.italic = True

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_comparison_table_slide(prs, title, headers, rows, notes=""):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title, size=28)
    rows_n, cols_n = len(rows) + 1, len(headers)
    left, top = Inches(0.5), Inches(1.4)
    width, height = Inches(12.3), Inches(0.6 * (rows_n + 1))
    table = slide.shapes.add_table(rows_n, cols_n,
                                   left, top, width, height).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for run in cell.text_frame.paragraphs[0].runs:
            run.font.bold = True
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        cell.fill.solid()
        cell.fill.fore_color.rgb = UTD_GREEN
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            for run in cell.text_frame.paragraphs[0].runs:
                run.font.size = Pt(14)
                run.font.color.rgb = DARK_TEXT
    slide.notes_slide.notes_text_frame.text = notes
    return slide


def add_results_slide(prs, title, image_path, headline_pairs, notes=""):
    """Hero results chart up top + headline number strip on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    _style_title(slide.shapes.title, size=28)

    if os.path.exists(image_path):
        slide.shapes.add_picture(image_path, Inches(0.3), Inches(1.4),
                                 width=Inches(8.5))

    box = slide.shapes.add_textbox(Inches(9.0), Inches(1.4),
                                   Inches(4.0), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = "Headline numbers (UTD live)"
    for run in tf.paragraphs[0].runs:
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = UTD_GREEN
    for label, value in headline_pairs:
        p = tf.add_paragraph()
        p.text = f"{label}"
        for run in p.runs:
            run.font.size = Pt(13)
            run.font.color.rgb = GREY_TEXT
        v = tf.add_paragraph()
        v.text = value
        for run in v.runs:
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = DARK_TEXT

    slide.notes_slide.notes_text_frame.text = notes
    return slide


def main():
    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ---- Slide 1: Title ----------------------------------------------------
    add_title_slide(
        prs,
        "UTD Parking Spot Occupancy Detector",
        "Turning the cameras you already have into a real-time spot-by-spot map",
        "Group 34  |  CS 6384 Computer Vision  |  "
        "Nikita Ramachandran, Sandeep Jammula, "
        "Praneeth Kumar Rachepalli, Eswardeep Pujala",
    )

    # ---- Slide 2: The hook ------------------------------------------------
    add_bullet_slide(
        prs,
        "The parking circle is a daily UTD problem",
        [
            "Commuter campus: 30k+ daily parkers, peak demand 9-11 AM.",
            "Today's structures show LOT-level totals (e.g. 'PS3: 47 free').",
            "Per-spot indicator lights are stuck-on / stuck-off across many spots.",
            "Drivers still climb every floor to actually find a space.",
            "Goal: tell a driver WHICH spots are open, before they enter the lot.",
        ],
        notes="Speaker: Nikita (~40 s). Open with the relatable scene - "
              "PS3 at 9:55 AM, sign says '47 free' but every spot light is "
              "lying. Land on: existing infrastructure measures the WRONG "
              "thing (lot total) or fails silently (per-spot sensors).",
    )

    # ---- Slide 3: What's wrong with existing solutions --------------------
    add_comparison_table_slide(
        prs,
        "Existing solutions vs. what drivers actually need",
        ["Approach", "What it gives you", "What's broken"],
        [
            ["Gate-loop counter",
             "Lot-level total only",
             "Driver still circles every floor"],
            ["Per-spot ultrasonic / IR / magnetometer sensor",
             "One bit per spot",
             "$50-150 per spot x thousands of spots; drift, weather, vandalism"],
            ["Patch classifier (CNRPark, PKLot)",
             "One CNN per spot crop",
             "Needs per-camera training data; brittle to new lighting"],
            ["Our approach",
             "Per-spot Occupied / Open in real time",
             "Uses existing camera; one human click pair per spot; no training"],
        ],
        notes="Speaker: Nikita (~45 s). Punchline: 'all three either ignore "
              "the driver's actual question or require expensive "
              "infrastructure / retraining for every camera. We use the "
              "cameras that already exist with zero per-camera retraining.'",
    )

    # ---- Slide 4: Methodology (pipeline visual) ---------------------------
    add_image_slide(
        prs,
        "How the model works  |  Frame -> YOLOv8 -> ROIs -> IoU > 0.5",
        os.path.join(FIGURES, "pipeline_visual.png"),
        bullets=[
            "1) Detector: YOLOv8n (Ultralytics, 6 MB), COCO weights, "
            "filtered to {car, truck, bus, motorcycle}.  "
            "2) Registration: click 2 corners per spot once -> rois.json.  "
            "3) Decision: per spot, if max IoU with any vehicle > 0.5 -> "
            "Occupied (red), else Open (green). No training.",
        ],
        notes="Speaker: Eswardeep (~60 s). Walk the audience through the "
              "four panels left-to-right. Stress: 'YOLOv8 already knows what "
              "cars look like from millions of COCO images, so we don't "
              "train anything. We just teach the system WHERE the spots "
              "are, once.'",
    )

    # ---- Slide 5: Setup + the two screenshots -----------------------------
    add_two_image_slide(
        prs,
        "Setup: register once, label ground truth once",
        os.path.join(FIGURES, "roi_labeling.png"),
        os.path.join(FIGURES, "gt_labeling.png"),
        left_caption="1. roi_picker.py - click two corners per spot, save rois.json",
        right_caption="2. label_gt.py - press o (occupied) / e (empty) per spot, 5 frames",
        notes="Speaker: Eswardeep -> Praneeth (~45 s). 41 s clip, 848x464, "
              "30 fps, 17 ROIs, 85 ground-truth judgments. That is the "
              "TOTAL manual effort - per camera, ever. Mention: this is "
              "the real screenshot from our tools, not a mockup of someone "
              "else's UI.",
    )

    # ---- Slide 6: Results (chart + headline numbers) ----------------------
    add_results_slide(
        prs,
        "Results on the live UTD recording",
        os.path.join(FIGURES, "results_chart.png"),
        [
            ("Accuracy",  "97.6 %"),
            ("Precision", "98.2 %"),
            ("Recall",    "98.2 %"),
            ("F1",        "98.2 %"),
            ("Inference (CPU, model only)", "17.3 FPS"),
            ("End-to-end FPS",              "15.7 FPS"),
            ("Errors",                      "2 / 85"),
        ],
        notes="Speaker: Praneeth (~60 s). Both columns of the chart are real "
              "measurements. Two errors total in 85 judgments - that's the "
              "number to memorise. 17 FPS on a laptop CPU = real-time. "
              "Then click into the demo on the next slide.",
    )

    # ---- Slide 7: Live demo ----------------------------------------------
    add_bullet_slide(
        prs,
        "Live demo",
        [
            "Red = Occupied,  Green = Open.",
            "HUD: live FPS  |  Occupied count (red)  |  Open count (green).",
            "Watch a spot flip as a car arrives or leaves.",
            "[Insert results/utd_demo.mp4 here in PowerPoint.]",
        ],
        notes="In PowerPoint: Insert > Video > This Device > "
              "results/utd_demo.mp4 > Start: Automatically. Trim to ~30 s "
              "if your section will overrun the 5:00 cap. Backup: keep "
              "data/frames/utd_demo_snapshot.jpg as a still on this slide.",
    )

    # ---- Slide 8: How we improved (recall iteration chart) ----------------
    add_image_slide(
        prs,
        "From 64% recall to 98% - by REDRAWING rectangles, not retraining",
        os.path.join(FIGURES, "precision_recall.png"),
        bullets=[
            "Round 1 (tight slivers along painted lines): 63.8% recall - "
            "ROIs traced visible asphalt, missed the actual car bodies.  "
            "Round 2 (car-sized boxes): 79.7%.  "
            "Round 3 (enlarged foreground rectangles to match perspective): "
            "98.2%. In a static-ROI system, the ROI IS the model.",
        ],
        notes="Speaker: Sandeep (~30 s). This is the clear improvement line "
              "the audience asks for. We never touched YOLO weights. "
              "Engineering rigor lives in the ROI design, not in the model.",
    )

    # ---- Slide 9: Failure analysis ---------------------------------------
    add_bullet_slide(
        prs,
        "Failure analysis (and how to fix each)",
        [
            "Top-down (~90 deg) cameras: COCO YOLO mis-classifies cars as "
            "ovens / microwaves -> mount cameras at 30-60 deg.",
            "Heavy occlusion in back rows: a foreground SUV can hide a "
            "back-row spot -> recover with a second camera or polygonal ROIs.",
            "Boundary-straddling vehicles: 1 FP + 1 FN on adjacent spots, "
            "a case where humans also disagree.",
            "Camera physical drift: every ROI shifts together -> periodic "
            "re-registration against painted-line fiducials.",
        ],
        notes="Speaker: Sandeep (~30 s). Be direct about the limits - the "
              "rubric explicitly rewards 'discuss and analyze failures.' "
              "Reference data/frames/carpark_diagnose_annotated.jpg as "
              "empirical proof of the top-down failure mode if asked.",
    )

    # ---- Slide 10: Conclusion --------------------------------------------
    add_bullet_slide(
        prs,
        "Conclusion: real-time, training-free, $0 per-spot hardware",
        [
            "97.6% accuracy and 98.2% F1 on a real UTD lot, at 17 FPS on a CPU.",
            "Uses cameras that already exist; new camera is ready in ~5 min "
            "(click ROIs once).",
            "Daily-life impact: a phone app could say 'Floor 3, North row, "
            "6 open spots' before a driver enters the structure.",
            "Future: auto-ROI from a short calibration video, multi-camera "
            "fusion, edge deployment on Jetson / Raspberry Pi.",
            "Thank you - questions?",
        ],
        notes="Speaker: Sandeep (~30 s). End BEFORE 5:00 to leave the full "
              "1 minute for Q&A. The 'daily-life impact' bullet is the "
              "punchline - hammer the spot-level vs lot-level distinction.",
    )

    # ---- Slide 11: Q&A backup --------------------------------------------
    add_bullet_slide(
        prs,
        "Q&A backup",
        [
            "Q: Why not train per-spot like CNRPark-EXT?  "
            "A: To avoid per-camera training data; detect-then-assign "
            "generalises to any new view in 2 min of clicks.",
            "Q: Why YOLOv8n vs v8m / v8l?  "
            "A: Real-time on CPU; bigger weights cost 5-10x latency for "
            "~1-2 mAP gain we don't need.",
            "Q: Why IoU 0.5?  "
            "A: PASCAL VOC convention + best F1 in our ablation over "
            "{0.3, 0.4, 0.5, 0.6}.",
            "Q: Could you fine-tune on UTD images?  "
            "A: Yes, but the point of this system is that you DON'T have to.",
        ],
        notes="Leave on screen during Q&A. Don't read aloud.",
    )

    prs.save(OUT)
    print(f"[OK] wrote {OUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
